"""Build a polished project presentation for the cow detection +
re-identification work.

Visual design:
  - 16:9 widescreen, dark gradient backgrounds
  - mint + magenta accent palette
  - hero stat cards, layered card panels
  - decorative dots / accent bars

Proof slides use ANNOTATED full images (cow_id boxes drawn) so the same
cow can be visually traced across multiple photos via its coloured box.

Requires: pip install python-pptx
"""

import csv
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


PROJECT = Path("/Users/yashwantyadav/Desktop/Cow_detection")
OUT_PPTX = PROJECT / "Cow_Detection_Project.pptx"
IMAGES_DIR = PROJECT / "Images"
ANNOTATED_DIR = PROJECT / "Cow_image_output" / "annotated"
COW_INDEX_CSV = PROJECT / "Cow_image_output" / "cow_index.csv"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ---- modern dark palette -----------------------------------------------------
BG_DARK = RGBColor(0x0B, 0x0F, 0x1F)
BG_MID = RGBColor(0x12, 0x17, 0x2E)
CARD = RGBColor(0x18, 0x1F, 0x39)
CARD_LIGHT = RGBColor(0x1E, 0x27, 0x46)
MINT = RGBColor(0x00, 0xD9, 0xB2)
PINK = RGBColor(0xFF, 0x6B, 0x9D)
AMBER = RGBColor(0xFB, 0xBF, 0x24)
TEXT = RGBColor(0xF5, 0xF7, 0xFA)
SUB = RGBColor(0x94, 0xA3, 0xB8)
DIM = RGBColor(0x6B, 0x72, 0x87)


# ---- helpers -----------------------------------------------------------------
def hide_outline(shape):
    shape.line.fill.background()


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    hide_outline(shape)


def add_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    fill(bg, BG_DARK)
    # subtle top-right corner glow stripe
    stripe = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(11.5), 0, Inches(2), Inches(0.5))
    fill(stripe, BG_MID)
    return slide


def add_text(slide, left, top, width, height, text,
             size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
             font="Helvetica Neue"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    return tf


def add_bullets(slide, left, top, width, height, items,
                size=18, color=TEXT, bullet_color=MINT,
                line_spacing=12, font="Helvetica Neue"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run_dot = p.add_run()
        run_dot.text = "▸  "
        run_dot.font.size = Pt(size)
        run_dot.font.bold = True
        run_dot.font.color.rgb = bullet_color
        run_dot.font.name = font
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = font
        p.space_after = Pt(line_spacing)


def add_header(slide, kicker, title, accent=MINT):
    # kicker (small uppercase label)
    add_text(slide, Inches(0.6), Inches(0.45), Inches(8), Inches(0.35),
             kicker.upper(), size=12, bold=True, color=accent)
    # main title
    add_text(slide, Inches(0.6), Inches(0.75), Inches(12), Inches(0.85),
             title, size=32, bold=True)
    # accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.55), Inches(0.5), Inches(0.06))
    fill(bar, accent)


def add_image(slide, path, left, top, width, height, caption=None,
              cap_color=SUB, border=False):
    if not Path(path).exists():
        ph = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    left, top, width, height)
        fill(ph, CARD)
        add_text(slide, left, top + height // 2 - Inches(0.2),
                 width, Inches(0.4),
                 f"missing: {Path(path).name}", size=12, color=DIM,
                 align=PP_ALIGN.CENTER)
        return
    pic = slide.shapes.add_picture(str(path), left, top, width=width,
                                   height=height)
    if border:
        pic.line.color.rgb = MINT
        pic.line.width = Pt(1.5)
    if caption:
        add_text(slide, left, top + height + Inches(0.05),
                 width, Inches(0.4),
                 caption, size=12, color=cap_color, align=PP_ALIGN.CENTER)


def stat_card(slide, left, top, width, height, value, label,
              accent=MINT):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD
    card.line.color.rgb = accent
    card.line.width = Pt(1.5)
    # value
    add_text(slide, left, top + Inches(0.25), width, Inches(1.0),
             value, size=44, bold=True, color=accent, align=PP_ALIGN.CENTER)
    # label
    add_text(slide, left, top + height - Inches(0.7), width, Inches(0.5),
             label.upper(), size=11, bold=True, color=SUB,
             align=PP_ALIGN.CENTER)


def panel(slide, left, top, width, height, color=CARD, accent=None):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    if accent:
        rect.line.color.rgb = accent
        rect.line.width = Pt(1.5)
    else:
        hide_outline(rect)
    return rect


def deco_dots(slide, left, top, count=6, color=MINT, size=8, gap=14):
    for i in range(count):
        d = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            left + i * Inches(gap / 72), top,
            Inches(size / 72), Inches(size / 72))
        fill(d, color)


# ---- slide builders ----------------------------------------------------------
def slide_title(prs):
    s = add_slide(prs)
    # left accent bar
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             0, Inches(2.5), Inches(0.15), Inches(2.5))
    fill(bar, MINT)
    # decoration dots
    deco_dots(s, Inches(0.6), Inches(2.0), count=5, gap=18)
    add_text(s, Inches(0.6), Inches(2.4), Inches(11), Inches(0.45),
             "COMPUTER VISION  •  LIVESTOCK", size=14, bold=True, color=MINT)
    add_text(s, Inches(0.6), Inches(2.9), Inches(12), Inches(1.5),
             "Cow Detection, Tracking", size=56, bold=True)
    add_text(s, Inches(0.6), Inches(3.85), Inches(12), Inches(1.5),
             "& Re-Identification", size=56, bold=True, color=MINT)
    add_text(s, Inches(0.6), Inches(5.1), Inches(11), Inches(0.6),
             "A computer-vision pipeline for automated livestock monitoring,",
             size=18, color=SUB)
    add_text(s, Inches(0.6), Inches(5.4), Inches(11), Inches(0.6),
             "individual identification, and herd analytics.",
             size=18, color=SUB)
    add_text(s, Inches(0.6), Inches(6.5), Inches(11), Inches(0.4),
             "YASHWANT YADAV   •   2026", size=11, bold=True, color=DIM)


def slide_problem(prs):
    s = add_slide(prs)
    add_header(s, "The challenge", "Why automated cow monitoring matters")
    # big number panel left
    panel(s, Inches(0.6), Inches(2.0), Inches(4.6), Inches(4.6), CARD, MINT)
    add_text(s, Inches(0.6), Inches(2.3), Inches(4.6), Inches(1.5),
             "Hours", size=20, bold=True, color=SUB, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), Inches(2.6), Inches(4.6), Inches(1.8),
             "of unwatched",
             size=22, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), Inches(3.2), Inches(4.6), Inches(2),
             "CCTV per day",
             size=22, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), Inches(4.6), Inches(4.6), Inches(1.3),
             "24/7", size=68, bold=True, color=MINT, align=PP_ALIGN.CENTER)
    # right column bullets
    add_bullets(s, Inches(5.8), Inches(2.0), Inches(7.2), Inches(5),
                [
                    "Manual headcounts are slow, error-prone, and cannot scale to large herds.",
                    "Existing CCTV records hours of footage every day — no one watches it.",
                    "Re-identifying individual cows is critical for health, breeding records, and loss detection.",
                    "Off-the-shelf models trained on drone imagery fail on ground-level CCTV (domain mismatch).",
                    "Trackers double-count the same cow when it leaves and re-enters frame, inflating counts.",
                ], size=18, line_spacing=14)


def slide_tracks(prs):
    s = add_slide(prs)
    add_header(s, "Our approach",
               "Two tools for two kinds of footage")
    # left card: video
    panel(s, Inches(0.6), Inches(1.95), Inches(5.95), Inches(5.1),
          CARD, MINT)
    add_text(s, Inches(0.85), Inches(2.1), Inches(1.5), Inches(0.4),
             "TOOL 1", size=11, bold=True, color=MINT)
    add_text(s, Inches(0.85), Inches(2.4), Inches(5.5), Inches(0.6),
             "For CCTV Videos", size=26, bold=True)
    add_text(s, Inches(0.85), Inches(3.0), Inches(5.5), Inches(0.4),
             "Counts cows in continuous footage", size=14, color=SUB)
    add_bullets(s, Inches(0.85), Inches(3.55), Inches(5.5), Inches(3.4),
                [
                    "Reads the camera's video file",
                    "Spots every cow in every frame",
                    "Follows each cow as it moves",
                    "Avoids counting the same cow twice",
                    "Outputs a video with cow numbers on screen",
                ], size=15, line_spacing=10)

    # right card: photos
    panel(s, Inches(6.8), Inches(1.95), Inches(5.95), Inches(5.1),
          CARD, PINK)
    add_text(s, Inches(7.05), Inches(2.1), Inches(1.5), Inches(0.4),
             "TOOL 2", size=11, bold=True, color=PINK)
    add_text(s, Inches(7.05), Inches(2.4), Inches(5.5), Inches(0.6),
             "For Photo Albums", size=26, bold=True)
    add_text(s, Inches(7.05), Inches(3.0), Inches(5.5), Inches(0.4),
             "Finds the same cow across different pictures",
             size=14, color=SUB)
    add_bullets(s, Inches(7.05), Inches(3.55), Inches(5.5), Inches(3.4),
                [
                    "Spots cows in each photo",
                    "Looks at the body, coat colour, and shape",
                    "Compares every cow to every other cow",
                    "Groups matching cows under one ID",
                    "Skips cows that are half-cut at the edges",
                ], size=15, line_spacing=10, bullet_color=PINK)


def slide_video_pipeline(prs):
    s = add_slide(prs)
    add_header(s, "How the video tool works",
               "From hours of footage to a clean cow count")
    panel(s, Inches(0.6), Inches(1.95), Inches(12.1), Inches(2.2), CARD)
    add_text(s, Inches(0.85), Inches(2.05), Inches(8), Inches(0.4),
             "STEP 1", size=11, bold=True, color=MINT)
    add_text(s, Inches(0.85), Inches(2.35), Inches(8), Inches(0.6),
             "Watch every frame and follow each cow", size=22, bold=True)
    add_bullets(s, Inches(0.85), Inches(3.05), Inches(11.7), Inches(1),
                [
                    "An AI vision model spots every cow in the picture",
                    "A tracker links the same cow across frames, even through brief occlusions",
                    "Each cow gets a number that stays with it as it walks around",
                ], size=14, line_spacing=4)

    panel(s, Inches(0.6), Inches(4.4), Inches(12.1), Inches(2.6), CARD)
    add_text(s, Inches(0.85), Inches(4.5), Inches(8), Inches(0.4),
             "STEP 2", size=11, bold=True, color=PINK)
    add_text(s, Inches(0.85), Inches(4.8), Inches(8), Inches(0.6),
             "Clean up the count", size=22, bold=True)
    add_bullets(s, Inches(0.85), Inches(5.5), Inches(11.7), Inches(1.5),
                [
                    "If a cow leaves the frame and returns, the system recognises it as the same animal",
                    "Two cows visible at the same time can never share a number",
                    "Final output is a single number per cow plus a clean labelled video",
                ], size=14, line_spacing=4, bullet_color=PINK)


def slide_video_results(prs):
    s = add_slide(prs)
    add_header(s, "Video tool  ·  results",
               "Frames from a 5-minute cow-shed CCTV clip")
    add_image(s, PROJECT / "merged_t30.jpg",
              Inches(0.5), Inches(2.0), Inches(6.1), Inches(3.4),
              caption="30 seconds in  ·  each cow has a number",
              border=True)
    add_image(s, PROJECT / "merged_t90.jpg",
              Inches(6.75), Inches(2.0), Inches(6.1), Inches(3.4),
              caption="90 seconds in  ·  the same numbers are still on the right cows",
              border=True)
    stat_card(s, Inches(0.5), Inches(5.85), Inches(4.0), Inches(1.3),
              "341 → 102", "before vs after de-duplication", MINT)
    stat_card(s, Inches(4.7), Inches(5.85), Inches(4.0), Inches(1.3),
              "5 min", "of cctv processed", MINT)
    stat_card(s, Inches(8.9), Inches(5.85), Inches(3.8), Inches(1.3),
              "≈ 3×", "fewer fake duplicates", PINK)


def slide_dup_bug(prs):
    s = add_slide(prs)
    add_header(s, "A problem we had to solve",
               "Two different cows ending up with the same number")
    panel(s, Inches(0.6), Inches(2.0), Inches(5.95), Inches(4.9),
          CARD, PINK)
    add_text(s, Inches(0.85), Inches(2.15), Inches(5.5), Inches(0.4),
             "THE PROBLEM", size=11, bold=True, color=PINK)
    add_text(s, Inches(0.85), Inches(2.45), Inches(5.5), Inches(0.6),
             "Mistaken identity", size=22, bold=True)
    add_bullets(s, Inches(0.85), Inches(3.2), Inches(5.5), Inches(3.6),
                [
                    "The system tried to merge cows that looked similar",
                    "Sometimes it merged a chain that connected two completely different animals",
                    "Two cows visible at the same time would share the same number",
                    "This inflated the headcount and broke per-cow tracking",
                ], size=14, line_spacing=8, bullet_color=PINK)

    panel(s, Inches(6.8), Inches(2.0), Inches(5.95), Inches(4.9),
          CARD, MINT)
    add_text(s, Inches(7.05), Inches(2.15), Inches(5.5), Inches(0.4),
             "THE FIX", size=11, bold=True, color=MINT)
    add_text(s, Inches(7.05), Inches(2.45), Inches(5.5), Inches(0.6),
             "Stricter matching rules", size=22, bold=True)
    add_bullets(s, Inches(7.05), Inches(3.2), Inches(5.5), Inches(3.6),
                [
                    "Never merge two cows that are on screen at the same time",
                    "Require a stronger appearance match before merging",
                    "Require the two sightings to be in nearby parts of the frame",
                    "Require a small time gap to avoid merging on weak evidence",
                ], size=14, line_spacing=8)


def slide_reid_pipeline(prs):
    s = add_slide(prs)
    add_header(s, "How the photo tool works",
               "From a folder of pictures to grouped cow identities")
    steps = [
        ("01", "Find the cows",
         "Detect every cow\nin every image"),
        ("02", "Keep only the clear ones",
         "Drop cows that are tiny\nor half-cut at the edge"),
        ("03", "Take a fingerprint",
         "Capture the body's\nappearance for each cow"),
        ("04", "Add coat colour",
         "Record the colour pattern\nas a second signature"),
        ("05", "Compare every pair",
         "Cows with similar\nfingerprints get grouped"),
        ("06", "Label the results",
         "Same group  =  same cow\nin a results spreadsheet"),
    ]
    w = Inches(4.0); h = Inches(2.3)
    coords = [
        (Inches(0.5), Inches(1.95)),
        (Inches(4.65), Inches(1.95)),
        (Inches(8.8), Inches(1.95)),
        (Inches(0.5), Inches(4.55)),
        (Inches(4.65), Inches(4.55)),
        (Inches(8.8), Inches(4.55)),
    ]
    for (num, title, body), (l, t) in zip(steps, coords):
        panel(s, l, t, w, h, CARD, MINT if int(num) <= 3 else PINK)
        add_text(s, l + Inches(0.25), t + Inches(0.15),
                 Inches(1.0), Inches(0.5),
                 num, size=26, bold=True,
                 color=MINT if int(num) <= 3 else PINK)
        add_text(s, l + Inches(0.25), t + Inches(0.65),
                 w - Inches(0.5), Inches(0.5),
                 title, size=18, bold=True)
        add_text(s, l + Inches(0.25), t + Inches(1.1),
                 w - Inches(0.5), Inches(1.1),
                 body, size=13, color=SUB)


def slide_feature_stack(prs):
    s = add_slide(prs)
    add_header(s, "What makes each cow unique",
               "Three things the system looks at when comparing cows")
    blocks = [
        ("Overall look", "Main signal",
         "The shape of the body, the texture of the coat, the layout of patches and spots.\n\n"
         "This is the strongest hint that two photos show the same animal.",
         MINT),
        ("Coat colour", "Backup signal",
         "How much white, black, brown, or red is in the cow's coat — and where.\n\n"
         "Works even when the lighting changes between photos.",
         AMBER),
        ("Body shape", "Tiebreaker",
         "How tall versus wide the cow is in the picture.\n\n"
         "Helps decide close calls when the first two signals are not enough.",
         PINK),
    ]
    cols = [(Inches(0.5), Inches(4.0)),
            (Inches(4.7), Inches(4.0)),
            (Inches(8.9), Inches(4.0))]
    for (title, role, desc, col), (l, w) in zip(blocks, cols):
        panel(s, l, Inches(1.95), w, Inches(5.1), CARD, col)
        add_text(s, l + Inches(0.3), Inches(2.1),
                 w - Inches(0.6), Inches(0.5),
                 role.upper(), size=12, bold=True, color=col)
        add_text(s, l + Inches(0.3), Inches(2.45),
                 w - Inches(0.6), Inches(0.6),
                 title, size=24, bold=True)
        add_text(s, l + Inches(0.3), Inches(3.4),
                 w - Inches(0.6), Inches(3),
                 desc, size=15, color=TEXT)
    add_text(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.5),
             "All three together  =  a unique fingerprint for that cow.",
             size=15, bold=True, color=MINT, align=PP_ALIGN.CENTER)


def slide_reid_results(prs):
    s = add_slide(prs)
    add_header(s, "Photo tool  ·  results",
               "Run on 58 farm images")
    stat_card(s, Inches(0.6), Inches(2.0), Inches(2.95), Inches(2.0),
              "58", "photos analysed")
    stat_card(s, Inches(3.65), Inches(2.0), Inches(2.95), Inches(2.0),
              "~530", "cow sightings")
    stat_card(s, Inches(6.7), Inches(2.0), Inches(2.95), Inches(2.0),
              "169", "distinct cows")
    stat_card(s, Inches(9.75), Inches(2.0), Inches(2.95), Inches(2.0),
              "20", "seen in more than one photo", PINK)

    panel(s, Inches(0.6), Inches(4.4), Inches(12.1), Inches(2.6),
          CARD, MINT)
    add_text(s, Inches(0.85), Inches(4.55), Inches(8), Inches(0.4),
             "WHY WE SKIP HALF-VISIBLE COWS", size=11, bold=True, color=MINT)
    add_text(s, Inches(0.85), Inches(4.85), Inches(11), Inches(0.6),
             "A cow that is half-cut at the edge can't be matched reliably.",
             size=18, bold=True)
    add_bullets(s, Inches(0.85), Inches(5.55), Inches(11.7), Inches(1.5),
                [
                    "Without the filter: 309 cows — many of them were the same animal counted twice",
                    "With the filter:    169 cows — cleaner, more honest count",
                    "We only fingerprint cows whose full body is visible in the photo",
                ], size=14, line_spacing=6)


def slide_side_by_side(prs):
    pairs = [
        ("WhatsApp Image 2026-05-01 at 09.13.17 (1).jpeg",
         "Same scene — first the raw photo, then the photo with cow numbers added."),
        ("WhatsApp Image 2026-05-01 at 09.13.11 (1).jpeg",
         "Cow #7 is first spotted here — and the system later finds it in three more photos."),
        ("WhatsApp Image 2026-05-01 at 09.13.22 (1).jpeg",
         "Herd shot — several cows here are matched to other photos in the set."),
    ]
    for name, caption in pairs:
        s = add_slide(prs)
        add_header(s, "Before & after", "What the tool produces")
        add_image(s, IMAGES_DIR / name,
                  Inches(0.5), Inches(2.0), Inches(6.1), Inches(4.4),
                  caption="ORIGINAL PHOTO", border=False)
        add_image(s, ANNOTATED_DIR / name,
                  Inches(6.75), Inches(2.0), Inches(6.1), Inches(4.4),
                  caption="WITH COW NUMBERS  ·  same colour = same cow elsewhere",
                  border=True)
        add_text(s, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.5),
                 caption, size=14, color=SUB, align=PP_ALIGN.CENTER)


def slide_same_cow(prs, cow_id, image_names, label_subtitle):
    """Same-cow proof using full annotated images (not crops)."""
    s = add_slide(prs)
    add_header(s, "Re-identification proof",
               f"cow_{cow_id:03d}  ·  {label_subtitle}")
    n = len(image_names)
    if n <= 2:
        w, h = Inches(6.1), Inches(4.6)
        positions = [(Inches(0.5), Inches(2.0)),
                     (Inches(6.75), Inches(2.0))]
    elif n == 3:
        w, h = Inches(4.0), Inches(3.0)
        positions = [(Inches(0.55), Inches(2.0)),
                     (Inches(4.65), Inches(2.0)),
                     (Inches(8.75), Inches(2.0))]
    elif n == 4:
        w, h = Inches(6.1), Inches(2.25)
        positions = [(Inches(0.5), Inches(2.0)),
                     (Inches(6.75), Inches(2.0)),
                     (Inches(0.5), Inches(4.45)),
                     (Inches(6.75), Inches(4.45))]
    else:  # 5
        w, h = Inches(4.0), Inches(2.25)
        positions = [(Inches(0.55), Inches(2.0)),
                     (Inches(4.65), Inches(2.0)),
                     (Inches(8.75), Inches(2.0)),
                     (Inches(2.6), Inches(4.45)),
                     (Inches(6.7), Inches(4.45))]
    for name, (l, t) in zip(image_names, positions):
        add_image(s, ANNOTATED_DIR / name, l, t, w, h,
                  caption=name.split(" at ")[-1].replace(".jpeg", ""),
                  border=True)
    add_text(s, Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.45),
             "Boxes with the same colour & label = same cow re-identified across images.",
             size=13, color=SUB, align=PP_ALIGN.CENTER)


def slide_accuracy(prs):
    """Honest verification numbers from cow_index.csv manual checks."""
    s = add_slide(prs)
    add_header(s, "Manual verification",
               "Visual audit of every re-identified cow")
    correct = 0; incorrect = 0; details = []
    if COW_INDEX_CSV.exists():
        with COW_INDEX_CSV.open() as f:
            for row in csv.DictReader(f):
                v = row.get("Manual_Check_Result", "").strip().lower()
                if v == "correct":
                    correct += 1
                elif v == "incorrect":
                    incorrect += 1
                details.append((int(row["cow_id"]),
                                int(row["num_images"]), v))
    total = correct + incorrect
    pct = (correct / total * 100) if total else 0
    stat_card(s, Inches(0.6), Inches(2.0), Inches(3.9), Inches(2.2),
              f"{pct:.0f}%", "overall precision", MINT)
    stat_card(s, Inches(4.7), Inches(2.0), Inches(3.9), Inches(2.2),
              str(correct), "correct matches", MINT)
    stat_card(s, Inches(8.8), Inches(2.0), Inches(3.9), Inches(2.2),
              str(incorrect), "incorrect matches", PINK)
    # top-confidence stat
    high_conf = [d for d in details if d[1] >= 3]
    high_correct = sum(1 for d in high_conf if d[2] == "correct")
    panel(s, Inches(0.6), Inches(4.5), Inches(12.1), Inches(2.4),
          CARD, MINT)
    add_text(s, Inches(0.85), Inches(4.65), Inches(8), Inches(0.4),
             "WHAT THE NUMBERS SHOW", size=11, bold=True, color=MINT)
    add_bullets(s, Inches(0.85), Inches(5.05), Inches(11.7), Inches(1.8),
                [
                    f"We checked all {total} matches by eye  ·  {correct} were genuinely the same cow.",
                    f"For cows seen in 3 or more photos, the match was right "
                    f"{(high_correct/len(high_conf)*100) if high_conf else 0:.0f}% of the time.",
                    "The more often a cow appears, the more confident the result.",
                    "Next: collect more cow photos so the tool can learn farm-specific patterns.",
                ], size=14, line_spacing=6)


def slide_applications(prs):
    s = add_slide(prs)
    add_header(s, "Real-world impact",
               "Problems this pipeline can solve on a working farm")
    apps = [
        ("Automated headcount",
         "Replace daily manual counts with continuous CCTV-driven counting."),
        ("Theft & loss detection",
         "Alert when a previously seen cow stops appearing in footage."),
        ("Individual health monitoring",
         "Flag cows whose appearance frequency drops — early illness signal."),
        ("Breeding & lineage records",
         "Verifiable photo trail of every animal over time."),
        ("Insurance & audit trail",
         "Proof-of-life images per insured animal."),
        ("Multi-source flexibility",
         "Same pipeline works on phone, drone, and fixed-camera input."),
    ]
    coords = [
        (Inches(0.6), Inches(2.0)),
        (Inches(4.7), Inches(2.0)),
        (Inches(8.8), Inches(2.0)),
        (Inches(0.6), Inches(4.55)),
        (Inches(4.7), Inches(4.55)),
        (Inches(8.8), Inches(4.55)),
    ]
    for (title, desc), (l, t) in zip(apps, coords):
        panel(s, l, t, Inches(4.0), Inches(2.3), CARD, MINT)
        add_text(s, l + Inches(0.25), t + Inches(0.2),
                 Inches(3.5), Inches(0.6),
                 title, size=18, bold=True, color=MINT)
        add_text(s, l + Inches(0.25), t + Inches(0.95),
                 Inches(3.5), Inches(1.3),
                 desc, size=13, color=TEXT)


def slide_tech(prs):
    s = add_slide(prs)
    add_header(s, "Tech stack",
               "Open-source, end-to-end on a single Mac")
    techs = [
        ("Detection", "Ultralytics YOLOv8m  ·  PyTorch  ·  MPS / CUDA / CPU"),
        ("Tracking", "BoT-SORT — Kalman + IoU + Global Motion Compensation"),
        ("Re-ID Embedding", "torchvision ResNet50 (ImageNet weights)"),
        ("Clustering", "scikit-learn AgglomerativeClustering  ·  cosine, average linkage"),
        ("Imaging", "OpenCV  —  HSV histograms, crops, drawing"),
        ("Video I/O", "ffmpeg H.264 piped from Python (small file output)"),
        ("Custom", "Proprietary .cpv NVR reader  +  Union-Find post-hoc merge"),
    ]
    for i, (label, value) in enumerate(techs):
        t = Inches(1.9 + i * 0.7)
        panel(s, Inches(0.6), t, Inches(12.1), Inches(0.6), CARD)
        add_text(s, Inches(0.85), t + Inches(0.13),
                 Inches(3.0), Inches(0.4),
                 label.upper(), size=12, bold=True, color=MINT)
        add_text(s, Inches(3.8), t + Inches(0.13),
                 Inches(8.5), Inches(0.4),
                 value, size=14)


def slide_deliverables(prs):
    s = add_slide(prs)
    add_header(s, "Deliverables",
               "Code, data, and outputs in this repository")
    items = [
        ("view_cpv.py", "NVR .cpv viewer / converter to MP4"),
        ("count_cows.py", "Track A — video tracking + dedup pipeline (CLI)"),
        ("cow_reid.py", "Track B — image re-identification pipeline (CLI)"),
        ("annotate_matches.py", "Draws cow_id boxes on every image with a matched cow"),
        ("build_cow_index.py", "Pivots summary.csv to one row per re-identified cow"),
        ("botsort_cows.yaml", "BoT-SORT tracker config tuned for static CCTV"),
        ("Cow_image_output/", "summary.csv · cow_index.csv · per-cow folders · annotated/ · montages"),
    ]
    for i, (name, desc) in enumerate(items):
        t = Inches(1.9 + i * 0.7)
        panel(s, Inches(0.6), t, Inches(12.1), Inches(0.6), CARD)
        # bullet dot
        d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), t + Inches(0.22),
                               Inches(0.16), Inches(0.16))
        fill(d, MINT)
        add_text(s, Inches(1.3), t + Inches(0.13),
                 Inches(4.0), Inches(0.4),
                 name, size=14, bold=True, color=MINT,
                 font="Menlo")
        add_text(s, Inches(5.0), t + Inches(0.13),
                 Inches(7.5), Inches(0.4),
                 desc, size=13, color=TEXT)


def slide_thanks(prs):
    s = add_slide(prs)
    deco_dots(s, Inches(5.2), Inches(2.3), count=8, gap=14)
    add_text(s, Inches(0.6), Inches(2.7), Inches(12.1), Inches(1.2),
             "Thank you", size=72, bold=True, align=PP_ALIGN.CENTER,
             color=MINT)
    add_text(s, Inches(0.6), Inches(4.0), Inches(12.1), Inches(0.6),
             "github.com/Deadshot1831/Cow_detection",
             size=20, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), Inches(4.55), Inches(12.1), Inches(0.5),
             "yadavyeshwant6166@gmail.com",
             size=14, color=SUB, align=PP_ALIGN.CENTER)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(6.0), Inches(5.5),
                             Inches(1.3), Inches(0.05))
    fill(bar, MINT)


# ---- main --------------------------------------------------------------------
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_problem(prs)
    slide_tracks(prs)
    slide_video_pipeline(prs)
    slide_video_results(prs)
    slide_dup_bug(prs)
    slide_reid_pipeline(prs)
    slide_feature_stack(prs)
    slide_reid_results(prs)
    slide_side_by_side(prs)

    # same-cow proof slides using FULL annotated images
    proofs = [
        (5, ["WhatsApp Image 2026-05-01 at 09.13.17 (1).jpeg",
             "WhatsApp Image 2026-05-01 at 09.13.17 (2).jpeg",
             "WhatsApp Image 2026-05-01 at 09.13.17.jpeg",
             "WhatsApp Image 2026-05-01 at 09.13.18.jpeg",
             "WhatsApp Image 2026-05-01 at 09.13.19.jpeg"],
         "spotted across 5 images (09.13.17 → 09.13.19)"),
        (7, ["WhatsApp Image 2026-05-01 at 09.13.11 (1).jpeg",
             "WhatsApp Image 2026-05-01 at 09.13.11 (2).jpeg",
             "WhatsApp Image 2026-05-01 at 09.13.12 (1).jpeg",
             "WhatsApp Image 2026-05-01 at 09.13.12.jpeg"],
         "spotted across 4 images (09.13.11 → 09.13.12)"),
        (1, ["WhatsApp Image 2026-05-01 at 09.13.11 (1).jpeg",
             "WhatsApp Image 2026-05-01 at 09.13.18 (1).jpeg",
             "WhatsApp Image 2026-05-01 at 09.13.18 (2).jpeg"],
         "spotted across 3 images (09.13.11 → 09.13.18)"),
        (8, ["WhatsApp Image 2026-05-01 at 09.13.04 (2).jpeg",
             "WhatsApp Image 2026-05-01 at 09.13.14 (2).jpeg",
             "WhatsApp Image 2026-05-01 at 09.13.15 (1).jpeg"],
         "spotted across 3 images (09.13.04 → 09.13.15)"),
    ]
    for cid, imgs, sub in proofs:
        slide_same_cow(prs, cid, imgs, sub)

    slide_accuracy(prs)
    slide_applications(prs)
    slide_thanks(prs)

    prs.save(str(OUT_PPTX))
    print(f"wrote {OUT_PPTX}")
    print(f"slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
